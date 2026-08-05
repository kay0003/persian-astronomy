"""
generate_presentation.py
Generates Persian_Astronomy_Presentation.pptx — a polished PowerPoint about
the history of Persian astronomy, designed for Polish high-school students.

Requirements: pip install -r requirements.txt
"""

import io
import os
import tempfile
import warnings
import logging
import requests

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree

try:
    import cairosvg
    HAS_CAIROSVG = True
except ImportError:
    HAS_CAIROSVG = False

logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
DARK_NAVY  = RGBColor(10, 14, 38)
GOLD       = RGBColor(232, 199, 125)
WHITE      = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(220, 220, 220)
GOLD_FADED = RGBColor(180, 155, 95)

# ---------------------------------------------------------------------------
# Slide dimensions  (16:9 widescreen)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Image catalogue
# ---------------------------------------------------------------------------
IMAGES = {
    1:  ("milky_way.jpg",    "https://upload.wikimedia.org/wikipedia/commons/1/14/Milky_Way_Arch_over_Lut_Desert%2C_Iran_by_Amirreza_Kamkar.jpg",
                              "Milky_Way_Arch_over_Lut_Desert,_Iran_by_Amirreza_Kamkar.jpg"),
    2:  ("achaemenid.jpg",   "https://upload.wikimedia.org/wikipedia/commons/4/44/Map_of_the_Achaemenid_Empire.jpg",
                              "Map_of_the_Achaemenid_Empire.jpg"),
    3:  ("fire_temple.jpg",  "https://upload.wikimedia.org/wikipedia/commons/1/13/Zoroastrian_Fire_Temple_in_Yazd.JPG",
                              "Zoroastrian_Fire_Temple_in_Yazd.JPG"),
    4:  ("faravahar.svg",    "https://upload.wikimedia.org/wikipedia/commons/5/5b/Faravahar.svg",
                              "Faravahar.svg"),
    5:  ("nowruz.jpg",       "https://upload.wikimedia.org/wikipedia/commons/8/82/A_traditional_Nowruz_%22Haft_Sin%22_table_at_the_Smithsonian_Institution%27s_Freer-Sackler_Museum_in_Washington_displays_seven_items_representing_renewal_and_new_life.jpg",
                              "A_traditional_Nowruz_\"Haft_Sin\"_table_at_the_Smithsonian_Institution's_Freer-Sackler_Museum_in_Washington_displays_seven_items_representing_renewal_and_new_life.jpg"),
    6:  ("khayyam.jpg",      "https://upload.wikimedia.org/wikipedia/commons/4/47/Hakim_Omar_Khayam_-_panoramio.jpg",
                              "Hakim_Omar_Khayam_-_panoramio.jpg"),
    7:  ("tusi_portrait.jpg","https://commons.wikimedia.org/wiki/Special:FilePath/Nasir_al-Din_al-Tusi_portrait.jpg",
                              "Nasir_al-Din_al-Tusi_portrait.jpg"),
    8:  ("tusi_couple.jpg",  "https://upload.wikimedia.org/wikipedia/commons/8/81/Tusi_couple.jpg",
                              "Tusi_couple.jpg"),
    9:  ("maragheh.jpg",     "https://upload.wikimedia.org/wikipedia/commons/9/95/Maragheh_Observatory_02.JPG",
                              "Maragheh_Observatory_02.JPG"),
    10: ("biruni.jpg",       "https://commons.wikimedia.org/wiki/Special:FilePath/Al-Biruni_Portrait.jpg",
                              "Al-Biruni_Portrait.jpg"),
    11: ("astrolabe.jpg",    "https://upload.wikimedia.org/wikipedia/commons/8/8e/Astrolabe-Persian-18C.jpg",
                              "Astrolabe-Persian-18C.jpg"),
    12: ("persepolis.jpg",   "https://commons.wikimedia.org/wiki/Special:FilePath/General_view_of_the_ruins_of_Persepolis.jpg",
                              "General_view_of_the_ruins_of_Persepolis.jpg"),
    13: ("chogha.jpg",       "https://upload.wikimedia.org/wikipedia/commons/5/5c/Choghazanbil2.jpg",
                              "Choghazanbil2.jpg"),
}

_img_cache: dict[int, str | None] = {}

def _download(url: str, dest: str) -> bool:
    """Download url to dest file. Returns True on success."""
    headers = {"User-Agent": "Mozilla/5.0 (compatible; PptxBot/1.0)"}
    try:
        r = requests.get(url, headers=headers, timeout=30, allow_redirects=True)
        r.raise_for_status()
        with open(dest, "wb") as f:
            f.write(r.content)
        return True
    except Exception as e:
        log.warning("Download failed %s — %s", url, e)
        return False


def get_image(img_id: int, tmpdir: str) -> str | None:
    """Return local path to image, downloading if needed. Returns None on failure."""
    if img_id in _img_cache:
        return _img_cache[img_id]

    filename, primary_url, fallback_name = IMAGES[img_id]
    dest = os.path.join(tmpdir, filename)

    # Try primary URL
    if _download(primary_url, dest):
        pass
    else:
        # Try Special:FilePath fallback
        fallback_url = f"https://commons.wikimedia.org/wiki/Special:FilePath/{requests.utils.quote(fallback_name)}"
        if not _download(fallback_url, dest):
            log.warning("Skipping image %d — all download attempts failed.", img_id)
            _img_cache[img_id] = None
            return None

    # Handle SVG (image 4)
    if filename.endswith(".svg"):
        png_dest = dest.replace(".svg", ".png")
        if HAS_CAIROSVG:
            try:
                cairosvg.svg2png(url=dest, write_to=png_dest, output_width=600)
                _img_cache[img_id] = png_dest
                return png_dest
            except Exception as e:
                log.warning("SVG→PNG conversion failed: %s", e)
        log.warning("Skipping SVG image %d (cairosvg unavailable or failed).", img_id)
        _img_cache[img_id] = None
        return None

    _img_cache[img_id] = dest
    return dest


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_name="Calibri", font_size=18, bold=False, italic=False,
                 color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_title(slide, title_text, top=Inches(0.25)):
    return add_text_box(
        slide, title_text,
        left=Inches(0.4), top=top, width=Inches(12.5), height=Inches(0.85),
        font_name="Georgia", font_size=36, bold=True, color=GOLD,
        align=PP_ALIGN.LEFT
    )


def add_bullets(slide, bullets: list[str],
                left=Inches(0.4), top=Inches(1.2),
                width=Inches(7.0), height=Inches(5.5)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.name = "Calibri"
        run.font.size = Pt(22)
        run.font.color.rgb = LIGHT_GRAY
    return txBox


def add_image_right(slide, img_path: str | None,
                    left=Inches(7.7), top=Inches(1.1),
                    width=Inches(5.3), height=Inches(5.7),
                    white_bg=False):
    """Add image with gold border on the right side of the slide."""
    if img_path is None:
        return

    # Optional white background box (for line-art images)
    if white_bg:
        bg_box = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left - Inches(0.05), top - Inches(0.05),
            width + Inches(0.1), height + Inches(0.1)
        )
        bg_box.fill.solid()
        bg_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg_box.line.color.rgb = GOLD
        bg_box.line.width = Pt(2)

    try:
        pic = slide.shapes.add_picture(img_path, left, top, width, height)
        # Add gold border via XML
        sp_tree = slide.shapes._spTree
        # Move picture to front
        sp_tree.append(pic._element)

        # Gold border on picture
        spPr = pic._element.find(qn("p:spPr"))
        if spPr is None:
            spPr = etree.SubElement(pic._element, qn("p:spPr"))
        ln = etree.SubElement(spPr, qn("a:ln"))
        ln.set("w", str(int(Pt(2.5))))
        solidFill = etree.SubElement(ln, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", str(GOLD))
    except Exception as e:
        log.warning("Could not add image: %s", e)


def add_centered_text_block(slide, text, top=Inches(2.0), width=Inches(9.0),
                             font_size=26, font_name="Calibri", color=LIGHT_GRAY,
                             italic=False):
    left = (SLIDE_W - width) / 2
    add_text_box(slide, text, left=left, top=top, width=width, height=Inches(3.0),
                 font_name=font_name, font_size=font_size, italic=italic,
                 color=color, align=PP_ALIGN.CENTER)


def add_horizontal_rule(slide, top=Inches(1.15)):
    """Thin gold line under the title."""
    from pptx.util import Emu
    line = slide.shapes.add_shape(1, Inches(0.4), top, Inches(12.5), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_title(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_NAVY)

    img = get_image(1, tmpdir)
    if img:
        # Full-bleed background image with dark overlay effect
        pic = slide.shapes.add_picture(img, 0, 0, SLIDE_W, SLIDE_H)
        # Darken with a semi-transparent overlay
        overlay = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(10, 14, 38)
        from pptx.dml.color import RGBColor as _RGB
        overlay.fill.fore_color.theme_color  # touch to force solid
        # Use transparency via XML
        sp_elem = overlay._element
        spPr = sp_elem.find(qn("p:spPr"))
        if spPr is not None:
            solidFill = spPr.find(f".//{qn('a:solidFill')}")
            if solidFill is not None:
                clr = solidFill.find(qn("a:srgbClr"))
                if clr is not None:
                    alpha = etree.SubElement(clr, qn("a:alpha"))
                    alpha.set("val", "75000")  # 75% opacity = 25% transparent
        overlay.line.fill.background()

    # Title
    add_text_box(slide, "Persian Astronomy",
                 left=Inches(1.0), top=Inches(2.0), width=Inches(11.3), height=Inches(1.4),
                 font_name="Georgia", font_size=60, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER)
    # Subtitle
    add_text_box(slide,
                 "Stars, Science & Civilization\n"
                 "— A Journey Through the Sky of Ancient & Medieval Persia —",
                 left=Inches(1.0), top=Inches(3.5), width=Inches(11.3), height=Inches(1.8),
                 font_name="Calibri", font_size=24, italic=True, color=LIGHT_GRAY,
                 align=PP_ALIGN.CENTER)


def slide_why_persia(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_title(slide, "Why Persia?")
    add_horizontal_rule(slide)
    add_bullets(slide, [
        "One of the oldest continuous civilizations (3000+ years)",
        "Crossroads of trade, religion, and science",
        "Astronomy tied to calendar, agriculture, and religion",
    ])
    add_image_right(slide, get_image(2, tmpdir))


def slide_zoroastrianism(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_title(slide, "Zoroastrianism & the Sky")
    add_horizontal_rule(slide)
    add_bullets(slide, [
        "Ancient Persians saw stars as sacred beings",
        "4 royal stars guarded the 4 directions — Tishtrya linked to Sirius",
        "Fire temples aligned with celestial events",
    ])
    add_image_right(slide, get_image(3, tmpdir))


def slide_faravahar(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_title(slide, "Sacred Symbol: Faravahar")
    add_horizontal_rule(slide)
    add_centered_text_block(slide,
        "Representing the spirit and Zoroastrian cosmic order",
        top=Inches(1.5), font_size=28, italic=True, color=LIGHT_GRAY)
    add_image_right(slide, get_image(4, tmpdir), white_bg=True)


def slide_calendar(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_title(slide, "The Persian Calendar")
    add_horizontal_rule(slide)
    add_bullets(slide, [
        "Solar Hijri calendar — one of the most accurate solar calendars ever devised",
        "Nowruz (New Year) precisely timed to the spring equinox — still used today!",
    ])
    add_image_right(slide, get_image(5, tmpdir))


def slide_khayyam(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_title(slide, "Omar Khayyam (1048–1131)")
    add_horizontal_rule(slide)
    add_bullets(slide, [
        "Poet AND mathematician-astronomer",
        "Reformed the Persian (Jalali) calendar — more accurate than the Gregorian calendar",
        "Worked at the Isfahan observatory",
    ])
    add_image_right(slide, get_image(6, tmpdir))


def slide_tusi(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_title(slide, "Nasir al-Din al-Tusi (1201–1274)")
    add_horizontal_rule(slide)
    add_bullets(slide, [
        "Founded the great Maragheh Observatory",
        "Invented the \"Tusi Couple\" — later echoed in Copernicus' work",
        "Compiled major star catalogs",
    ])
    img_path = get_image(7, tmpdir)
    add_image_right(slide, img_path, top=Inches(1.1), height=Inches(5.0))
    if img_path:
        add_text_box(slide,
            "(Artistic/imagined portrait — no real historical depiction survives)",
            left=Inches(7.7), top=Inches(6.2), width=Inches(5.3), height=Inches(0.5),
            font_size=12, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_tusi_couple(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_title(slide, "The Tusi Couple")
    add_horizontal_rule(slide)
    add_centered_text_block(slide,
        "A geometric device explaining planetary motion\n"
        "— an idea later echoed by Copernicus",
        top=Inches(1.5), font_size=26, italic=True, color=LIGHT_GRAY)
    add_image_right(slide, get_image(8, tmpdir), white_bg=True)


def slide_maragheh(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_title(slide, "The Maragheh Observatory")
    add_horizontal_rule(slide)
    add_bullets(slide, [
        "Built in 1259 CE, scholars from Persia, China & beyond",
        "Library of 400,000+ books",
        "Inspired later observatories across Asia and possibly Europe",
    ])
    add_image_right(slide, get_image(9, tmpdir))


def slide_biruni(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_title(slide, "Al-Biruni (973–1048)")
    add_horizontal_rule(slide)
    add_bullets(slide, [
        "Calculated Earth's radius with striking accuracy",
        "Discussed Earth possibly orbiting the Sun — centuries before Copernicus",
        "Wrote on latitude, longitude, eclipses, star positions",
    ])
    add_image_right(slide, get_image(10, tmpdir))


def slide_astrolabe(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_title(slide, "The Astrolabe")
    add_horizontal_rule(slide)
    add_centered_text_block(slide,
        "A beautiful brass instrument refined by Persian scientists\n"
        "— used for navigation, timekeeping, and mapping stars.",
        top=Inches(1.5), font_size=26, italic=False, color=LIGHT_GRAY)
    add_image_right(slide, get_image(11, tmpdir))


def slide_persepolis(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_title(slide, "Sky-Aligned Architecture: Persepolis")
    add_horizontal_rule(slide)
    add_centered_text_block(slide,
        "Columns and gates aligned with solstice sunrise",
        top=Inches(1.5), font_size=28, italic=True, color=LIGHT_GRAY)
    add_image_right(slide, get_image(12, tmpdir))


def slide_chogha(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_title(slide, "Chogha Zanbil Ziggurat")
    add_horizontal_rule(slide)
    add_centered_text_block(slide,
        "One of the few surviving ziggurats outside Mesopotamia\n"
        "— UNESCO World Heritage Site",
        top=Inches(1.5), font_size=26, color=LIGHT_GRAY)
    add_image_right(slide, get_image(13, tmpdir))


def slide_legacy(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_title(slide, "Legacy Today")
    add_horizontal_rule(slide)
    add_centered_text_block(slide,
        "Persian astronomy terms passed into Arabic → Latin → English.\n"
        "Iran continues astronomy research today.",
        top=Inches(1.5), font_size=26, color=LIGHT_GRAY)
    add_image_right(slide, get_image(1, tmpdir))


def slide_section_divider(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)

    # Decorative star autoshape (5-pointed star)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches as I
    star = slide.shapes.add_shape(
        92,  # fivePointedStar
        I(6.0), I(0.3), I(1.3), I(1.0)
    )
    star.fill.solid()
    star.fill.fore_color.rgb = GOLD
    star.line.fill.background()

    add_text_box(slide, "Let's Write in Persian! ✍️",
                 left=Inches(1.0), top=Inches(2.2), width=Inches(11.3), height=Inches(1.4),
                 font_name="Georgia", font_size=48, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, "We'll learn 4 simple words, letter by letter",
                 left=Inches(1.5), top=Inches(3.8), width=Inches(10.3), height=Inches(0.9),
                 font_name="Calibri", font_size=26, italic=True, color=LIGHT_GRAY,
                 align=PP_ALIGN.CENTER)


def _set_cell_text(cell, text, font_name="Calibri", font_size=18,
                   bold=False, color=WHITE, align=PP_ALIGN.CENTER, persian=False):
    tf = cell.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    p.clear()
    run = p.add_run()
    run.text = text
    run.font.name = "Arial" if persian else font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def slide_letters(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_title(slide, "Letters We Need")
    add_horizontal_rule(slide)

    rows_data = [
        ("س", "seen", '"s"'),
        ("ت", "te", '"t"'),
        ("ا", "alef", '"a"'),
        ("ر", "re", '"r"'),
        ("ه", "he", '"e / h"'),
        ("آ", "alef madde", '"â"'),
        ("م", "mim", '"m"'),
        ("ن", "noon", '"n"'),
        ("خ", "khe", '"kh"'),
        ("و", "vav", '"o / v"'),
        ("ش", "shin", '"sh"'),
        ("ی", "ye", '"i / y"'),
        ("د", "dal", '"d"'),
    ]
    n_rows = len(rows_data) + 1  # +1 header

    from pptx.util import Inches as I, Pt
    tbl_left = I(0.5)
    tbl_top = I(1.2)
    tbl_width = I(8.5)
    tbl_height = I(5.8)

    table = slide.shapes.add_table(n_rows, 3, tbl_left, tbl_top, tbl_width, tbl_height).table
    col_widths = [I(1.8), I(3.5), I(3.2)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header
    headers = ["Letter", "Name", "Sound"]
    for col, hdr in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(50, 50, 90)
        _set_cell_text(cell, hdr, bold=True, font_size=18, color=GOLD)

    for row_i, (letter, name, sound) in enumerate(rows_data, start=1):
        row_color = RGBColor(20, 25, 60) if row_i % 2 == 0 else RGBColor(15, 20, 50)

        cell0 = table.cell(row_i, 0)
        cell0.fill.solid()
        cell0.fill.fore_color.rgb = row_color
        _set_cell_text(cell0, letter, font_name="Arial", font_size=34,
                       bold=True, color=GOLD, persian=True)

        cell1 = table.cell(row_i, 1)
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = row_color
        _set_cell_text(cell1, name, font_size=20, color=LIGHT_GRAY)

        cell2 = table.cell(row_i, 2)
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = row_color
        _set_cell_text(cell2, sound, font_size=20, color=LIGHT_GRAY)

    # Note below table
    add_text_box(slide, "Remember: Persian is written right-to-left!",
                 left=I(0.5), top=I(7.05), width=I(8.5), height=I(0.35),
                 font_size=15, italic=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_vocab(prs, tmpdir, persian_word, pronunciation, english, polish):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_horizontal_rule(slide, top=Inches(0.95))

    # Large Persian word
    add_text_box(slide, persian_word,
                 left=Inches(0.5), top=Inches(1.1), width=Inches(12.3), height=Inches(2.5),
                 font_name="Arial", font_size=72, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER)
    # Pronunciation
    add_text_box(slide, pronunciation,
                 left=Inches(0.5), top=Inches(3.5), width=Inches(12.3), height=Inches(0.7),
                 font_name="Calibri", font_size=28, italic=True, color=LIGHT_GRAY,
                 align=PP_ALIGN.CENTER)

    # Small 2-column table
    from pptx.util import Inches as I
    tbl_left = I(4.0)
    tbl_top = I(4.4)
    tbl_w = I(5.3)
    tbl_h = I(1.2)
    table = slide.shapes.add_table(2, 2, tbl_left, tbl_top, tbl_w, tbl_h).table
    table.columns[0].width = I(2.65)
    table.columns[1].width = I(2.65)

    for col_i, (hdr, val) in enumerate([("English", english), ("Polish", polish)]):
        hc = table.cell(0, col_i)
        hc.fill.solid()
        hc.fill.fore_color.rgb = RGBColor(50, 50, 90)
        _set_cell_text(hc, hdr, bold=True, color=GOLD, font_size=16)

        vc = table.cell(1, col_i)
        vc.fill.solid()
        vc.fill.fore_color.rgb = RGBColor(20, 25, 60)
        _set_cell_text(vc, val, color=LIGHT_GRAY, font_size=20)


def slide_practice(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_title(slide, "Your Turn! ✏️")
    add_horizontal_rule(slide)

    add_text_box(slide,
                 "Trace and practice writing all 4 words, right to left.",
                 left=Inches(0.5), top=Inches(1.2), width=Inches(12.3), height=Inches(0.6),
                 font_size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Dashed-border rectangle
    rect = slide.shapes.add_shape(1, Inches(0.7), Inches(2.0), Inches(11.9), Inches(4.8))
    rect.fill.background()
    ln = rect._element.find(qn("p:spPr")).find(qn("a:ln"))
    if ln is None:
        spPr = rect._element.find(qn("p:spPr"))
        ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(Pt(1.5))))
    solidFill = etree.SubElement(ln, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", str(GOLD))
    prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", "dash")

    # Four words spaced apart inside the rectangle
    words = ["ستاره", "آسمان", "ماه", "خورشید"]
    positions = [Inches(0.9), Inches(3.7), Inches(6.5), Inches(9.3)]
    for word, left in zip(words, positions):
        add_text_box(slide, word,
                     left=left, top=Inches(3.0), width=Inches(2.5), height=Inches(1.5),
                     font_name="Arial", font_size=42, bold=True,
                     color=GOLD_FADED, align=PP_ALIGN.CENTER)


def slide_thankyou(prs, tmpdir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)

    img = get_image(1, tmpdir)
    if img:
        pic = slide.shapes.add_picture(img, 0, 0, SLIDE_W, SLIDE_H)
        overlay = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(10, 14, 38)
        sp_elem = overlay._element
        spPr = sp_elem.find(qn("p:spPr"))
        if spPr is not None:
            solidFill = spPr.find(f".//{qn('a:solidFill')}")
            if solidFill is not None:
                clr = solidFill.find(qn("a:srgbClr"))
                if clr is not None:
                    alpha = etree.SubElement(clr, qn("a:alpha"))
                    alpha.set("val", "75000")
        overlay.line.fill.background()

    add_text_box(slide, "متشکرم",
                 left=Inches(1.0), top=Inches(1.5), width=Inches(11.3), height=Inches(1.8),
                 font_name="Arial", font_size=72, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, "motshakeram",
                 left=Inches(1.0), top=Inches(3.2), width=Inches(11.3), height=Inches(0.7),
                 font_name="Calibri", font_size=28, italic=True, color=LIGHT_GRAY,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, "Thank you  •  dziękuję",
                 left=Inches(1.0), top=Inches(4.0), width=Inches(11.3), height=Inches(0.7),
                 font_name="Calibri", font_size=26, color=LIGHT_GRAY,
                 align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build_presentation(output_path="Persian_Astronomy_Presentation.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    with tempfile.TemporaryDirectory() as tmpdir:
        log.info("=== Downloading images ===")
        slide_title(prs, tmpdir)                   # 1
        slide_why_persia(prs, tmpdir)              # 2
        slide_zoroastrianism(prs, tmpdir)          # 3
        slide_faravahar(prs, tmpdir)               # 4
        slide_calendar(prs, tmpdir)                # 5
        slide_khayyam(prs, tmpdir)                 # 6
        slide_tusi(prs, tmpdir)                    # 7
        slide_tusi_couple(prs, tmpdir)             # 8
        slide_maragheh(prs, tmpdir)                # 9
        slide_biruni(prs, tmpdir)                  # 10
        slide_astrolabe(prs, tmpdir)               # 11
        slide_persepolis(prs, tmpdir)              # 12
        slide_chogha(prs, tmpdir)                  # 13
        slide_legacy(prs, tmpdir)                  # 14
        slide_section_divider(prs, tmpdir)         # 15
        slide_letters(prs, tmpdir)                 # 16
        # Vocabulary slides 17-20
        slide_vocab(prs, tmpdir, "ستاره", "setâreh", "star", "gwiazda")
        slide_vocab(prs, tmpdir, "آسمان", "âsemân", "sky", "niebo")
        slide_vocab(prs, tmpdir, "ماه", "mâh", "moon", "księżyc")
        slide_vocab(prs, tmpdir, "خورشید", "khorshid", "sun", "słońce")
        slide_practice(prs, tmpdir)               # 21
        slide_thankyou(prs, tmpdir)               # 22

    log.info("=== Saving %s ===", output_path)
    prs.save(output_path)
    log.info("Done! Saved to %s", output_path)


if __name__ == "__main__":
    build_presentation()
